"""
PowerPoint Hybrid Controller with Slide Deck Content Parser.
Supports both Desktop PowerPoint (COM / Process) and Web Presenter Canvas.
"""

import os
import time
import threading
import subprocess
import winreg
import ctypes
import pythoncom
import win32com.client
import win32gui
import win32con
import pyautogui
from pptx import Presentation

try:
    ctypes.windll.user32.AllowSetForegroundWindow(-1)
except Exception:
    pass


def find_powerpnt_exe():
    """Find the full executable path of POWERPNT.EXE on Windows."""
    try:
        key = winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\powerpnt.exe")
        val, _ = winreg.QueryValueEx(key, "")
        if val and os.path.exists(val):
            return val
    except Exception:
        pass

    candidates = [
        r"C:\Program Files\Microsoft Office\Root\Office16\POWERPNT.EXE",
        r"C:\Program Files (x86)\Microsoft Office\Root\Office16\POWERPNT.EXE",
        r"C:\Program Files\Microsoft Office\Office16\POWERPNT.EXE",
        r"C:\Program Files (x86)\Microsoft Office\Office16\POWERPNT.EXE",
        r"C:\Program Files\Microsoft Office\Office15\POWERPNT.EXE",
    ]
    for path in candidates:
        if os.path.exists(path):
            return path

    return "POWERPNT.EXE"


def bring_ppt_to_front():
    """Bring PowerPoint or SlideShow window to active foreground."""
    try:
        def callback(hwnd, extra):
            if win32gui.IsWindowVisible(hwnd):
                title = win32gui.GetWindowText(hwnd)
                if ("PowerPoint" in title or "Slide Show" in title) and not ("Assistant" in title or "Chrome" in title or "Edge" in title or "Code" in title):
                    ctypes.windll.user32.ShowWindow(hwnd, win32con.SW_RESTORE)
                    ctypes.windll.user32.SetForegroundWindow(hwnd)
        win32gui.EnumWindows(callback, None)
    except Exception as e:
        print(f"Foreground activation note: {e}")


class PPTController:
    def __init__(self, default_folder=None):
        self.default_folder = default_folder or os.path.abspath("./presentations")
        os.makedirs(self.default_folder, exist_ok=True)
        self.active_presentation_path = None
        self.current_slide = 1
        self.total_slides = 1
        self.ppt_exe = find_powerpnt_exe()
        self._lock = threading.Lock()

    def _ensure_com(self):
        """Initialize COM for current thread."""
        pythoncom.CoInitialize()

    def get_app(self, create=False):
        """Get or optionally start PowerPoint application."""
        self._ensure_com()
        try:
            app = win32com.client.GetActiveObject("PowerPoint.Application")
            return app
        except Exception:
            if create:
                try:
                    app = win32com.client.Dispatch("PowerPoint.Application")
                    app.Visible = True
                    return app
                except Exception:
                    return None
            return None

    def list_presentations(self, folder_path=None):
        """List all PowerPoint files in target directory."""
        target_dir = folder_path or self.default_folder
        if not os.path.exists(target_dir):
            return []

        ppt_files = []
        for file in os.listdir(target_dir):
            if file.lower().endswith((".pptx", ".ppt", ".ppsx", ".pptm")) and not file.startswith("~$"):
                full_path = os.path.abspath(os.path.join(target_dir, file))
                ppt_files.append({
                    "name": file,
                    "path": full_path,
                    "size_kb": round(os.path.getsize(full_path) / 1024, 1),
                    "modified": os.path.getmtime(full_path),
                    "is_active": (self.active_presentation_path == full_path)
                })
        return sorted(ppt_files, key=lambda x: x["name"])

    def parse_deck_slides(self, file_path=None):
        """Extract structured slide data using python-pptx."""
        path = file_path or self.active_presentation_path
        if not path or not os.path.exists(path):
            files = self.list_presentations()
            if files:
                path = files[0]["path"]
            else:
                return []

        try:
            prs = Presentation(path)
            slides_data = []
            for idx, slide in enumerate(prs.slides, 1):
                title_text = "Slide " + str(idx)
                subtitle_text = ""
                bullets = []

                if slide.shapes.title and slide.shapes.title.text.strip():
                    title_text = slide.shapes.title.text.strip()

                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        for p_idx, p in enumerate(shape.text_frame.paragraphs):
                            text = p.text.strip()
                            if text:
                                if idx == 1 and not subtitle_text and text != title_text:
                                    subtitle_text = text
                                else:
                                    bullets.append(text)

                slides_data.append({
                    "index": idx,
                    "title": title_text,
                    "subtitle": subtitle_text,
                    "content": bullets,
                    "is_title_slide": (idx == 1 or "welcome" in title_text.lower() or "presentation" in title_text.lower())
                })
            self.total_slides = max(1, len(slides_data))
            return slides_data
        except Exception as e:
            print(f"Error parsing deck slides: {e}")
            return []

    def open_presentation(self, file_path_or_name=None, start_show=True):
        """
        Open presentation in both Desktop PowerPoint and Web Stage.
        """
        with self._lock:
            self._ensure_com()
            files = self.list_presentations()

            if not files:
                raise FileNotFoundError(f"No PowerPoint files found in '{self.default_folder}'.")

            clean_target = (file_path_or_name or "").strip().lower()
            if clean_target in ["", "ppt", "the ppt", "presentation", "the presentation", "deck", "the deck", "default", "first", "it"]:
                clean_target = ""

            target_path = None
            if not clean_target:
                target_path = files[0]["path"]
            elif os.path.isabs(file_path_or_name) and os.path.exists(file_path_or_name):
                target_path = file_path_or_name
            else:
                match = None
                for f in files:
                    if clean_target in f["name"].lower():
                        match = f["path"]
                        break
                if match:
                    target_path = match
                else:
                    direct = os.path.join(self.default_folder, file_path_or_name)
                    if os.path.exists(direct):
                        target_path = direct
                    else:
                        target_path = files[0]["path"]

            self.active_presentation_path = target_path
            self.current_slide = 1

            # Extract slides for web viewer
            slides = self.parse_deck_slides(target_path)
            self.total_slides = len(slides) if slides else 1

            # Try launching desktop PowerPoint in parallel
            try:
                if start_show:
                    subprocess.Popen([self.ppt_exe, "/s", target_path], shell=False)
                else:
                    subprocess.Popen([self.ppt_exe, target_path], shell=False)
            except Exception:
                try:
                    os.startfile(target_path)
                except Exception:
                    pass

            return {
                "status": "opened",
                "file": os.path.basename(target_path),
                "path": target_path,
                "slideshow_running": True,
                "current_slide": self.current_slide,
                "total_slides": self.total_slides,
                "slides": slides
            }

    def start_slideshow(self):
        """Start fullscreen slideshow mode."""
        with self._lock:
            self._ensure_com()
            app = self.get_app(create=False)
            if app and app.Presentations.Count > 0:
                try:
                    app.ActivePresentation.SlideShowSettings.Run()
                    bring_ppt_to_front()
                except Exception:
                    pass
            return {"status": "slideshow_started", "current_slide": self.current_slide, "total_slides": self.total_slides}

    def stop_slideshow(self):
        """Exit slideshow."""
        with self._lock:
            self._ensure_com()
            app = self.get_app(create=False)
            if app:
                try:
                    if app.SlideShowWindows.Count > 0:
                        for i in range(1, app.SlideShowWindows.Count + 1):
                            app.SlideShowWindows(i).View.Exit()
                except Exception:
                    pass
            pyautogui.press("esc")
            return {"status": "slideshow_stopped"}

    def next_slide(self):
        """Advance to next slide."""
        with self._lock:
            self._ensure_com()
            if not self.active_presentation_path:
                return self.open_presentation(start_show=True)

            if self.current_slide < self.total_slides:
                self.current_slide += 1

            # Sync to Desktop PowerPoint
            app = self.get_app(create=False)
            if app:
                try:
                    if app.SlideShowWindows.Count > 0:
                        app.SlideShowWindows(1).View.Next()
                    elif app.Presentations.Count > 0:
                        app.ActiveWindow.View.GotoSlide(self.current_slide)
                except Exception:
                    pass
            else:
                pyautogui.press("right")

            return {"status": "success", "current_slide": self.current_slide, "total_slides": self.total_slides}

    def prev_slide(self):
        """Go to previous slide."""
        with self._lock:
            self._ensure_com()
            if not self.active_presentation_path:
                return self.open_presentation(start_show=True)

            if self.current_slide > 1:
                self.current_slide -= 1

            # Sync to Desktop PowerPoint
            app = self.get_app(create=False)
            if app:
                try:
                    if app.SlideShowWindows.Count > 0:
                        app.SlideShowWindows(1).View.Previous()
                    elif app.Presentations.Count > 0:
                        app.ActiveWindow.View.GotoSlide(self.current_slide)
                except Exception:
                    pass
            else:
                pyautogui.press("left")

            return {"status": "success", "current_slide": self.current_slide, "total_slides": self.total_slides}

    def goto_slide(self, slide_num: int):
        """Jump to specific slide."""
        with self._lock:
            self._ensure_com()
            if not self.active_presentation_path:
                self.open_presentation(start_show=True)

            target = max(1, min(int(slide_num), self.total_slides))
            self.current_slide = target

            # Sync to Desktop PowerPoint
            app = self.get_app(create=False)
            if app:
                try:
                    if app.SlideShowWindows.Count > 0:
                        app.SlideShowWindows(1).View.GotoSlide(target)
                    elif app.Presentations.Count > 0:
                        app.ActiveWindow.View.GotoSlide(target)
                except Exception:
                    pass

            return {"status": "success", "current_slide": self.current_slide, "total_slides": self.total_slides}

    def first_slide(self):
        return self.goto_slide(1)

    def last_slide(self):
        return self.goto_slide(self.total_slides)

    def blank_screen(self, color="black"):
        pyautogui.press("b" if color == "black" else "w")
        return {"status": "success", "screen_state": color}

    def close_presentation(self):
        with self._lock:
            self._ensure_com()
            app = self.get_app(create=False)
            if app and app.Presentations.Count > 0:
                try:
                    if app.SlideShowWindows.Count > 0:
                        app.SlideShowWindows(1).View.Exit()
                    app.ActivePresentation.Close()
                except Exception:
                    pass
            self.active_presentation_path = None
            self.current_slide = 1
            return {"status": "closed"}

    def get_status(self):
        self._ensure_com()
        app = self.get_app(create=False)
        pres_name = os.path.basename(self.active_presentation_path) if self.active_presentation_path else None
        
        # If COM has active window, sync numbers
        if app:
            try:
                if app.SlideShowWindows.Count > 0:
                    ss_win = app.SlideShowWindows(1)
                    pres_name = ss_win.Presentation.Name
                    self.total_slides = ss_win.Presentation.Slides.Count
                    self.current_slide = ss_win.View.Slide.SlideIndex
                elif app.Presentations.Count > 0:
                    pres_name = app.ActivePresentation.Name
                    self.total_slides = app.ActivePresentation.Slides.Count
                    self.current_slide = app.ActiveWindow.View.Slide.SlideIndex
            except Exception:
                pass

        return {
            "running": bool(self.active_presentation_path or app),
            "presentation_name": pres_name,
            "presentation_path": self.active_presentation_path,
            "current_slide": self.current_slide,
            "total_slides": self.total_slides,
            "in_slideshow": True if self.active_presentation_path else False,
            "folder": self.default_folder
        }
