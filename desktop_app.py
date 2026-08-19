"""
PowerPoint AI Presenter - Desktop GUI Application (Tkinter + Voice + Chat + Slide Stage)
"""

import os
import sys
import threading
import time
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
from pptx import Presentation

from backend.ppt_controller import PPTController
from backend.command_parser import parse_command

try:
    import speech_recognition as sr
    HAS_SR = True
except ImportError:
    HAS_SR = False

try:
    import pyttsx3
    HAS_TTS = True
except ImportError:
    HAS_TTS = False


class DesktopPresenterApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("PowerPoint AI Presenter (Voice & Chat Assistant)")
        self.geometry("1100x720")
        self.minsize(900, 600)
        self.configure(bg="#0b0f19")

        self.controller = PPTController(os.path.abspath("./presentations"))
        self.slides = []
        self.current_slide_idx = 1
        self.listening = False
        self.recognizer = sr.Recognizer() if HAS_SR else None

        self.setup_styles()
        self.build_ui()
        self.load_default_presentation()

    def setup_styles(self):
        style = ttk.Style(self)
        style.theme_use("clam")
        style.configure(".", background="#0b0f19", foreground="#f8fafc", font=("Segoe UI", 10))
        style.configure("TFrame", background="#0b0f19")
        style.configure("Card.TFrame", background="#151d30", relief="flat")
        style.configure("TLabel", background="#0b0f19", foreground="#f8fafc", font=("Segoe UI", 10))
        style.configure("Header.TLabel", background="#151d30", foreground="#ffffff", font=("Segoe UI", 14, "bold"))
        style.configure("Sub.TLabel", background="#151d30", foreground="#94a3b8", font=("Segoe UI", 9))
        style.configure("Accent.TButton", font=("Segoe UI", 10, "bold"), background="#2563eb", foreground="#ffffff")

    def build_ui(self):
        # 1. Top Header Bar
        top_bar = tk.Frame(self, bg="#111827", height=50, padx=16, pady=8)
        top_bar.pack(fill=tk.X, side=tk.TOP)

        title_lbl = tk.Label(top_bar, text="🎯 PPT Voice & Chat Presenter", font=("Segoe UI", 13, "bold"), fg="#ffffff", bg="#111827")
        title_lbl.pack(side=tk.LEFT)

        self.deck_status_lbl = tk.Label(top_bar, text="Deck: Loading...", font=("Segoe UI", 10, "bold"), fg="#60a5fa", bg="#111827")
        self.deck_status_lbl.pack(side=tk.LEFT, padx=20)

        self.slide_counter_lbl = tk.Label(top_bar, text="Slide: 1 / 1", font=("Segoe UI", 11, "bold"), fg="#34d399", bg="#111827")
        self.slide_counter_lbl.pack(side=tk.RIGHT)

        # 2. Main Content Split (Left: Stage + Controls, Right: Chat + Voice)
        main_split = tk.Frame(self, bg="#0b0f19", padx=12, pady=12)
        main_split.pack(fill=tk.BOTH, expand=True)

        # Left Column (Slide Stage + Remotes)
        left_col = tk.Frame(main_split, bg="#0b0f19")
        left_col.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(0, 8))

        # Slide Stage Canvas Box
        self.stage_frame = tk.Frame(left_col, bg="#1e1b4b", bd=2, relief=tk.GROOVE)
        self.stage_frame.pack(fill=tk.BOTH, expand=True, pady=(0, 10))

        self.slide_title_lbl = tk.Label(
            self.stage_frame, 
            text="Voice & Chat Controlled Presentation", 
            font=("Segoe UI", 20, "bold"), 
            fg="#ffffff", 
            bg="#1e1b4b",
            wraplength=600,
            justify=tk.CENTER
        )
        self.slide_title_lbl.pack(pady=(30, 10), padx=20)

        self.slide_sub_lbl = tk.Label(
            self.stage_frame, 
            text="Smart PowerPoint Assistant", 
            font=("Segoe UI", 13), 
            fg="#93c5fd", 
            bg="#1e1b4b",
            wraplength=600,
            justify=tk.CENTER
        )
        self.slide_sub_lbl.pack(pady=(0, 20), padx=20)

        self.slide_body_text = tk.Text(
            self.stage_frame, 
            bg="#1e1b4b", 
            fg="#e2e8f0", 
            font=("Segoe UI", 13),
            relief=tk.FLAT,
            wrap=tk.WORD,
            padx=30,
            height=8
        )
        self.slide_body_text.pack(fill=tk.BOTH, expand=True, padx=20, pady=(0, 20))
        self.slide_body_text.configure(state=tk.DISABLED)

        # Remote Control Toolbar
        remote_bar = tk.Frame(left_col, bg="#111827", padx=10, pady=8)
        remote_bar.pack(fill=tk.X)

        btn_first = tk.Button(remote_bar, text="⏮ First", command=self.first_slide, bg="#1f2937", fg="#ffffff", font=("Segoe UI", 9, "bold"), relief=tk.FLAT, padx=10, pady=4)
        btn_first.pack(side=tk.LEFT, padx=3)

        btn_prev = tk.Button(remote_bar, text="◀ Previous Slide", command=self.prev_slide, bg="#2563eb", fg="#ffffff", font=("Segoe UI", 10, "bold"), relief=tk.FLAT, padx=14, pady=4)
        btn_prev.pack(side=tk.LEFT, padx=4)

        btn_next = tk.Button(remote_bar, text="Next Slide ▶", command=self.next_slide, bg="#2563eb", fg="#ffffff", font=("Segoe UI", 10, "bold"), relief=tk.FLAT, padx=14, pady=4)
        btn_next.pack(side=tk.LEFT, padx=4)

        btn_last = tk.Button(remote_bar, text="Last ⏭", command=self.last_slide, bg="#1f2937", fg="#ffffff", font=("Segoe UI", 9, "bold"), relief=tk.FLAT, padx=10, pady=4)
        btn_last.pack(side=tk.LEFT, padx=3)

        btn_ppt = tk.Button(remote_bar, text="🎬 Launch PowerPoint App", command=self.launch_powerpoint_desktop, bg="#059669", fg="#ffffff", font=("Segoe UI", 9, "bold"), relief=tk.FLAT, padx=10, pady=4)
        btn_ppt.pack(side=tk.RIGHT, padx=4)

        # Right Column (Voice Hub + Chat Assistant)
        right_col = tk.Frame(main_split, bg="#111827", width=340, padx=10, pady=10)
        right_col.pack(side=tk.RIGHT, fill=tk.BOTH)

        # Voice Assistant Hub
        voice_box = tk.Frame(right_col, bg="#1f2937", padx=10, pady=10)
        voice_box.pack(fill=tk.X, pady=(0, 10))

        self.voice_btn = tk.Button(
            voice_box, 
            text="🎙 Click to Start Voice Listening", 
            command=self.toggle_voice, 
            bg="#dc2626", 
            fg="#ffffff", 
            font=("Segoe UI", 10, "bold"), 
            relief=tk.FLAT, 
            pady=8
        )
        self.voice_btn.pack(fill=tk.X)

        self.voice_status_lbl = tk.Label(voice_box, text="Voice: Ready (Say 'Next slide', 'Open PPT', etc.)", font=("Segoe UI", 8), fg="#94a3b8", bg="#1f2937")
        self.voice_status_lbl.pack(pady=(4, 0))

        # Quick Suggestions Chips
        chips_frame = tk.Frame(right_col, bg="#111827")
        chips_frame.pack(fill=tk.X, pady=(0, 8))

        tk.Label(chips_frame, text="Quick Commands:", font=("Segoe UI", 8, "bold"), fg="#94a3b8", bg="#111827").pack(anchor="w")

        chip_box = tk.Frame(chips_frame, bg="#111827")
        chip_box.pack(fill=tk.X, pady=2)

        for text in ["Open the PPT", "Next slide", "Previous slide", "Go to slide 3", "First slide"]:
            b = tk.Button(chip_box, text=text, command=lambda t=text: self.execute_text_command(t), bg="#1e293b", fg="#cbd5e1", font=("Segoe UI", 8), relief=tk.FLAT, padx=6, pady=2)
            b.pack(side=tk.LEFT, padx=2, pady=2)

        # Chat Log
        self.chat_log = tk.Text(right_col, bg="#0f172a", fg="#f8fafc", font=("Segoe UI", 9), relief=tk.FLAT, wrap=tk.WORD, state=tk.DISABLED)
        self.chat_log.pack(fill=tk.BOTH, expand=True, pady=(0, 8))

        # Chat Input Box
        input_frame = tk.Frame(right_col, bg="#111827")
        input_frame.pack(fill=tk.X)

        self.chat_entry = tk.Entry(input_frame, bg="#1e293b", fg="#ffffff", font=("Segoe UI", 10), insertbackground="#ffffff", relief=tk.FLAT)
        self.chat_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6), ipady=4)
        self.chat_entry.bind("<Return>", lambda e: self.on_chat_submit())

        send_btn = tk.Button(input_frame, text="Send", command=self.on_chat_submit, bg="#2563eb", fg="#ffffff", font=("Segoe UI", 9, "bold"), relief=tk.FLAT, padx=12)
        send_btn.pack(side=tk.RIGHT)

        self.log_chat("System", "Assistant started! Say 'Open PPT', 'Next slide', or type in chat.")

    def log_chat(self, sender, text):
        self.chat_log.configure(state=tk.NORMAL)
        self.chat_log.insert(tk.END, f"[{sender}] {text}\n")
        self.chat_log.see(tk.END)
        self.chat_log.configure(state=tk.DISABLED)

    def load_default_presentation(self):
        files = self.controller.list_presentations()
        if files:
            self.load_deck(files[0]["path"])
        else:
            self.deck_status_lbl.configure(text="No presentations found.")

    def load_deck(self, path):
        self.controller.active_presentation_path = path
        self.slides = self.controller.parse_deck_slides(path)
        self.current_slide_idx = 1
        name = os.path.basename(path)
        self.deck_status_lbl.configure(text=f"Deck: {name}")
        self.render_slide(1)
        self.log_chat("Assistant", f"Loaded presentation: {name} ({len(self.slides)} slides)")

    def render_slide(self, index):
        if not self.slides:
            return
        index = max(1, min(index, len(self.slides)))
        self.current_slide_idx = index
        slide = self.slides[index - 1]

        self.slide_title_lbl.configure(text=slide.get("title", f"Slide {index}"))
        sub = slide.get("subtitle", "")
        self.slide_sub_lbl.configure(text=sub)

        self.slide_body_text.configure(state=tk.NORMAL)
        self.slide_body_text.delete("1.0", tk.END)
        content = slide.get("content", [])
        if content:
            for item in content:
                self.slide_body_text.insert(tk.END, f"•  {item}\n\n")
        self.slide_body_text.configure(state=tk.DISABLED)

        self.slide_counter_lbl.configure(text=f"Slide: {index} / {len(self.slides)}")

    def next_slide(self):
        if self.current_slide_idx < len(self.slides):
            self.render_slide(self.current_slide_idx + 1)
        self.controller.next_slide()
        self.log_chat("Assistant", f"Moved to slide {self.current_slide_idx} of {len(self.slides)}")

    def prev_slide(self):
        if self.current_slide_idx > 1:
            self.render_slide(self.current_slide_idx - 1)
        self.controller.prev_slide()
        self.log_chat("Assistant", f"Moved to slide {self.current_slide_idx} of {len(self.slides)}")

    def first_slide(self):
        self.render_slide(1)
        self.controller.first_slide()
        self.log_chat("Assistant", "Navigated to first slide.")

    def last_slide(self):
        self.render_slide(len(self.slides))
        self.controller.last_slide()
        self.log_chat("Assistant", "Navigated to last slide.")

    def launch_powerpoint_desktop(self):
        self.log_chat("Assistant", "Launching native Microsoft PowerPoint window...")
        threading.Thread(target=lambda: self.controller.open_presentation(start_show=True), daemon=True).start()

    def on_chat_submit(self):
        text = self.chat_entry.get().strip()
        if not text:
            return
        self.chat_entry.delete(0, tk.END)
        self.log_chat("You", text)
        self.execute_text_command(text)

    def execute_text_command(self, text):
        parsed = parse_command(text)
        action = parsed.get("action")
        feedback = parsed.get("feedback", "")

        if action == "next":
            self.next_slide()
        elif action == "prev":
            self.prev_slide()
        elif action == "first":
            self.first_slide()
        elif action == "last":
            self.last_slide()
        elif action == "goto":
            slide = parsed.get("params", {}).get("slide", 1)
            self.render_slide(slide)
            self.controller.goto_slide(slide)
            self.log_chat("Assistant", f"Jumped to slide {slide}")
        elif action == "open":
            self.load_default_presentation()
            self.launch_powerpoint_desktop()
        elif action == "start_show":
            self.launch_powerpoint_desktop()
        else:
            self.log_chat("Assistant", feedback)

    def toggle_voice(self):
        if not HAS_SR:
            messagebox.showinfo("Voice", "SpeechRecognition is not installed. Use chat or remote buttons.")
            return

        if self.listening:
            self.listening = False
            self.voice_btn.configure(text="🎙 Click to Start Voice Listening", bg="#dc2626")
            self.voice_status_lbl.configure(text="Voice: Standby")
        else:
            self.listening = True
            self.voice_btn.configure(text="🔴 LISTENING... (Speak now)", bg="#16a34a")
            self.voice_status_lbl.configure(text="Listening for commands (Say 'Next slide', 'Open PPT')...")
            threading.Thread(target=self._voice_listen_loop, daemon=True).start()

    def _voice_listen_loop(self):
        mic = sr.Microphone()
        with mic as source:
            self.recognizer.adjust_for_ambient_noise(source, duration=0.6)

        while self.listening:
            try:
                with mic as source:
                    audio = self.recognizer.listen(source, timeout=3, phrase_time_limit=5)
                text = self.recognizer.recognize_google(audio)
                if text:
                    self.after(0, lambda t=text: self.log_chat("Voice", t))
                    self.after(0, lambda t=text: self.execute_text_command(t))
            except sr.WaitTimeoutError:
                pass
            except sr.UnknownValueError:
                pass
            except Exception as e:
                time.sleep(0.5)


if __name__ == "__main__":
    app = DesktopPresenterApp()
    app.mainloop()
