"""
Helper script to generate a rich sample PowerPoint presentation for testing.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_sample_presentation(target_path="presentations/Sample_AI_Presentation.pptx"):
    os.makedirs(os.path.dirname(os.path.abspath(target_path)), exist_ok=True)
    prs = Presentation()

    # Define color scheme
    NAVY = RGBColor(15, 23, 42)
    BLUE = RGBColor(37, 99, 235)
    LIGHT_GRAY = RGBColor(241, 245, 249)

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "Voice & Chat Controlled Presentation"
    subtitle.text = "Smart PowerPoint Assistant\nSay 'Next Slide' or use the Chat to Navigate"

    # Slide 2: Features Overview
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes2 = slide2.shapes
    shapes2.title.text = "Key Features & Capabilities"
    tf2 = shapes2.placeholders[1].text_frame
    tf2.text = "Hands-free voice recognition with real-time feedback."
    
    p = tf2.add_paragraph()
    p.text = "Natural chat commands ('open ppt', 'next slide', 'jump to slide 4')."
    p.level = 0
    
    p = tf2.add_paragraph()
    p.text = "Seamless COM integration with Microsoft PowerPoint."
    p.level = 0

    p = tf2.add_paragraph()
    p.text = "Support for any folder filled with PPT / PPTX presentations."
    p.level = 0

    # Slide 3: Voice Commands Table
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    shapes3 = slide3.shapes
    shapes3.title.text = "Voice Commands Reference"
    tf3 = shapes3.placeholders[1].text_frame
    tf3.text = "• 'Next Slide' / 'Advance' -> Advances to next slide"
    p = tf3.add_paragraph()
    p.text = "• 'Previous Slide' / 'Go Back' -> Returns to previous slide"
    p = tf3.add_paragraph()
    p.text = "• 'Go to slide [number]' -> Jumps directly to slide"
    p = tf3.add_paragraph()
    p.text = "• 'Start Presentation' -> Launches fullscreen slideshow"
    p = tf3.add_paragraph()
    p.text = "• 'Black Screen' / 'White Screen' -> Toggles screen blanking"
    p = tf3.add_paragraph()
    p.text = "• 'Exit Slideshow' -> Ends fullscreen mode"

    # Slide 4: Architecture
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    shapes4 = slide4.shapes
    shapes4.title.text = "System Architecture"
    tf4 = shapes4.placeholders[1].text_frame
    tf4.text = "1. Browser Frontend: Web Speech API & Real-time Chat HUD"
    p = tf4.add_paragraph()
    p.text = "2. Python Backend: FastAPI, WebSocket & Intent Parsing"
    p = tf4.add_paragraph()
    p.text = "3. Windows COM Engine: win32com PowerPoint automation"

    # Slide 5: Conclusion
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    shapes5 = slide5.shapes
    shapes5.title.text = "Thank You!"
    tf5 = shapes5.placeholders[1].text_frame
    tf5.text = "Presentation complete! Say 'Close PPT' or restart anytime."
    p = tf5.add_paragraph()
    p.text = "Enjoy your modern hands-free presentation experience!"

    prs.save(target_path)
    print(f"Sample presentation created at: {os.path.abspath(target_path)}")

if __name__ == "__main__":
    create_sample_presentation()
